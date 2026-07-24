"""report.py — turn a prescriptive plan into business deliverables.

`write_deliverables(plan, outdir="deliverables")` emits the artefacts a RevOps
team actually consumes:

    prescription.json         the full machine-readable plan
    optimization_workbook.xlsx an openpyxl workbook (Summary / Assortment /
                              Inventory / Pricing / Promo), styled headers
    executive_review.pdf      an ~8-slide matplotlib deck (uplift waterfall,
                              assortment before/after, inventory frontier,
                              price-move distribution, promo allocation, a
                              model-quality slide, recommended actions)
    executive_review.pptx     the same deck as PowerPoint (only if python-pptx
                              is importable — skipped gracefully otherwise)
    actions.csv               the actionable per-SKU list (reorder + reprice)

Only the plan dict is required; the inventory service-vs-cost frontier slide is
enriched from the SKU master when it is available, and degrades gracefully when
it is not. Charts share one palette and print one chart per page with numbers.

Author: Dimitres Kisimov.
"""
from __future__ import annotations

import csv
import json
from pathlib import Path

# ---- palette (shared across every chart) ---------------------------------- #
BLUE = "#2f6bff"
PINK = "#ea4b71"
GREEN = "#1d9e6f"
INK = "#1f2933"
GREY = "#9aa5b1"
AUTHOR = "Dimitres Kisimov"


# --------------------------------------------------------------------------- #
# helpers                                                                      #
# --------------------------------------------------------------------------- #
def _fmt_eur(v: float) -> str:
    return f"EUR {v:,.0f}"


def _carried_lines(plan: dict) -> list[dict]:
    return plan["plan"]["pricing"]["lines"]


def _sample_frontier(plan: dict):
    """(sku, frontier rows) for a carried SKU, or None if masters unavailable."""
    try:
        from .dataio import load
        from .optimize.inventory import service_cost_frontier
    except Exception:
        return None
    carried = set(plan["plan"]["assortment"]["carried"])
    skus = {s.sku: s for s in load()}
    # pick a carried SKU with real variability so the curve is legible
    cands = [skus[s] for s in carried if s in skus and skus[s].demand_std > 0]
    if not cands:
        return None
    s = max(cands, key=lambda x: x.avg_monthly_demand)
    return s.sku, service_cost_frontier(s)


# --------------------------------------------------------------------------- #
# 1. prescription.json + actions.csv                                          #
# --------------------------------------------------------------------------- #
def _write_json(plan: dict, outdir: Path) -> Path:
    p = outdir / "prescription.json"
    with p.open("w", encoding="utf-8") as f:
        json.dump(plan, f, indent=2)
    return p


def _write_actions_csv(plan: dict, outdir: Path) -> Path:
    """Per-SKU actionable list: reorder qty (order-up-to) + price change."""
    p = outdir / "actions.csv"
    price_by = {li["sku"]: li for li in _carried_lines(plan)}
    inv_by = {li["sku"]: li for li in plan["plan"]["inventory"]["lines"]}
    carried = plan["plan"]["assortment"]["carried"]
    with p.open("w", newline="", encoding="utf-8") as f:
        w = csv.writer(f)
        w.writerow(["sku", "action_carry", "order_up_to", "reorder_point",
                    "safety_stock", "current_price_eur", "recommended_price_eur",
                    "price_change_pct", "monthly_margin_uplift_eur",
                    "annual_margin_uplift_eur", "price_flag"])
        for sku in carried:
            pr = price_by.get(sku, {})
            iv = inv_by.get(sku, {})
            w.writerow([
                sku, "carry",
                iv.get("order_up_to", ""), iv.get("reorder_point", ""),
                iv.get("safety_stock", ""),
                pr.get("current_price", ""), pr.get("recommended_price", ""),
                pr.get("price_change_pct", ""),
                pr.get("margin_uplift_eur", ""),
                round(pr.get("margin_uplift_eur", 0.0) * 12, 2),
                pr.get("flag", ""),
            ])
    return p


# --------------------------------------------------------------------------- #
# 2. optimization_workbook.xlsx                                               #
# --------------------------------------------------------------------------- #
def _write_workbook(plan: dict, outdir: Path) -> Path | None:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter
    except Exception:
        return None

    hdr_fill = PatternFill("solid", fgColor="2F6BFF")
    hdr_font = Font(bold=True, color="FFFFFF")
    title_font = Font(bold=True, size=13, color="1F2933")

    wb = Workbook()

    def style_header(ws, ncols: int, row: int = 1) -> None:
        for c in range(1, ncols + 1):
            cell = ws.cell(row=row, column=c)
            cell.fill = hdr_fill
            cell.font = hdr_font
            cell.alignment = Alignment(horizontal="center")

    def autosize(ws) -> None:
        for col in ws.columns:
            width = max((len(str(c.value)) for c in col if c.value is not None),
                        default=8)
            ws.column_dimensions[get_column_letter(col[0].column)].width = \
                min(40, width + 2)

    # --- Summary -----------------------------------------------------------
    ws = wb.active
    ws.title = "Summary"
    up = plan["expected_uplift"]
    a = plan["plan"]["assortment"]
    fq = plan["model_quality"]["forecast"]
    rq = plan["model_quality"]["decline_risk"]
    ws["A1"] = "Revenue & Operations Optimization — Executive Summary"
    ws["A1"].font = title_font
    ws["A2"] = f"Author: {AUTHOR}"
    ws.append([])
    ws.append(["Expected annual uplift by lever", "EUR"])
    style_header(ws, 2, ws.max_row)
    ws.append(["Pricing (elasticity Lerner)", up["pricing_uplift_annual_eur"]])
    ws.append(["Promo (concave allocation)", up["promo_incremental_eur"]])
    ws.append(["Assortment (MILP vs greedy)", up["assortment_milp_vs_greedy_eur"]])
    ws.append(["Total expected uplift", up["total_expected_uplift_eur"]])
    ws.append([])
    ws.append(["KPI", "Value"])
    style_header(ws, 2, ws.max_row)
    ws.append(["SKUs carried", a["n_carried"]])
    ws.append(["SKUs dropped", a["n_dropped"]])
    ws.append(["Working capital used (EUR)", a["capital_used_eur"]])
    ws.append(["Working capital budget (EUR)", a["budget_eur"]])
    ws.append(["Assortment margin — MILP (EUR/yr)", a["optimized_annual_margin_eur"]])
    ws.append(["Assortment margin — greedy (EUR/yr)", a["greedy_annual_margin_eur"]])
    ws.append(["Baseline current annual margin (EUR)",
               plan["baseline_current_annual_margin_eur"]])
    ws.append(["Forecast MASE (vs seasonal-naive)",
               f"{fq['mase']} vs {fq['naive_mase']}"])
    ws.append(["Decline-risk ROC-AUC / PR-AUC",
               f"{rq['roc_auc']} / {rq['pr_auc']}"])
    autosize(ws)

    # --- Assortment --------------------------------------------------------
    ws = wb.create_sheet("Assortment")
    ws.append(["sku", "decision", "annual_margin_eur", "capital_at_cost_eur"])
    style_header(ws, 4)
    margin_by, capital_by = _assort_economics(plan)
    for sku in a["carried"]:
        ws.append([sku, "carry", margin_by.get(sku, ""), capital_by.get(sku, "")])
    for sku in a["dropped"]:
        ws.append([sku, "drop", margin_by.get(sku, ""), capital_by.get(sku, "")])
    autosize(ws)

    # --- Inventory ---------------------------------------------------------
    ws = wb.create_sheet("Inventory")
    cols = ["sku", "order_up_to", "safety_stock", "reorder_point",
            "critical_fractile", "expected_shortage_units", "expected_holding_eur"]
    ws.append(cols)
    style_header(ws, len(cols))
    for li in plan["plan"]["inventory"]["lines"]:
        ws.append([li[c] for c in cols])
    autosize(ws)

    # --- Pricing -----------------------------------------------------------
    ws = wb.create_sheet("Pricing")
    cols = ["sku", "current_price", "recommended_price", "price_change_pct",
            "elasticity", "flag", "expected_margin_now", "expected_margin_new",
            "margin_uplift_eur"]
    ws.append(cols)
    style_header(ws, len(cols))
    for li in _carried_lines(plan):
        ws.append([li[c] for c in cols])
    autosize(ws)

    # --- Promo -------------------------------------------------------------
    ws = wb.create_sheet("Promo")
    ws.append(["category", "allocation_eur", "share_pct"])
    style_header(ws, 3)
    pm = plan["plan"]["promo"]
    tot = sum(pm["allocation"].values()) or 1.0
    for cat, v in sorted(pm["allocation"].items(), key=lambda kv: -kv[1]):
        ws.append([cat, round(v, 2), round(100 * v / tot, 1)])
    ws.append([])
    ws.append(["Incremental margin (EUR)", pm["incremental_margin_eur"]])
    ws.append(["Uplift vs even split (EUR)", pm["uplift_vs_even_eur"]])
    ws.append(["Budget (EUR)", pm["budget_eur"]])
    autosize(ws)

    p = outdir / "optimization_workbook.xlsx"
    wb.save(p)
    return p


def _assort_economics(plan: dict) -> tuple[dict, dict]:
    """Per-SKU annual margin & capital-at-cost from the master (for the sheet)."""
    try:
        from .dataio import load
        skus = load()
        return ({s.sku: round(s.annual_margin_eur, 2) for s in skus},
                {s.sku: round(s.avg_inventory_cost_eur, 2) for s in skus})
    except Exception:
        return {}, {}


# --------------------------------------------------------------------------- #
# 3. executive_review.pdf (+ .pptx)                                           #
# --------------------------------------------------------------------------- #
def _slide_title(fig, plan: dict) -> None:
    fig.text(0.5, 0.62, "Revenue & Operations Optimization",
             ha="center", fontsize=26, fontweight="bold", color=INK)
    fig.text(0.5, 0.54, "Executive Review", ha="center", fontsize=20, color=BLUE)
    up = plan["expected_uplift"]
    fig.text(0.5, 0.40,
             f"Total expected uplift  {_fmt_eur(up['total_expected_uplift_eur'])} / year",
             ha="center", fontsize=16, color=GREEN)
    fig.text(0.5, 0.33,
             "AI demand/elasticity/decline models feeding a MILP + newsvendor "
             "+ pricing + promo optimization core",
             ha="center", fontsize=11, color=GREY)
    fig.text(0.5, 0.12, AUTHOR, ha="center", fontsize=12, color=INK)


def _slide_waterfall(fig, plan: dict) -> None:
    ax = fig.add_subplot(111)
    up = plan["expected_uplift"]
    base = plan["baseline_current_annual_margin_eur"]
    steps = [
        ("Baseline\nmargin", base, GREY),
        ("+ Pricing", up["pricing_uplift_annual_eur"], BLUE),
        ("+ Promo", up["promo_incremental_eur"], PINK),
        ("+ Assortment\n(MILP)", up["assortment_milp_vs_greedy_eur"], GREEN),
    ]
    running = 0.0
    for i, (_label, val, color) in enumerate(steps):
        if i == 0:
            ax.bar(i, val, color=color, edgecolor="white")
            running = val
        else:
            ax.bar(i, val, bottom=running, color=color, edgecolor="white")
            ax.plot([i - 1 + 0.4, i + 0.4], [running, running], color=GREY,
                    lw=0.8, ls="--")
            running += val
        ax.text(i, running + base * 0.01, _fmt_eur(val),
                ha="center", va="bottom", fontsize=9, color=INK)
    # final optimized total
    ax.bar(len(steps), running, color=INK, edgecolor="white")
    ax.text(len(steps), running + base * 0.01, _fmt_eur(running),
            ha="center", va="bottom", fontsize=9, fontweight="bold", color=INK)
    ax.set_xticks(range(len(steps) + 1))
    ax.set_xticklabels([s[0] for s in steps] + ["Optimized\ntotal"])
    ax.set_ylabel("Annual gross margin (EUR)")
    ax.set_title("Uplift waterfall — baseline to optimized", fontsize=14,
                 fontweight="bold", color=INK)
    ax.spines[["top", "right"]].set_visible(False)
    ax.yaxis.grid(True, color="#eef1f5")
    ax.set_axisbelow(True)


def _slide_assortment(fig, plan: dict) -> None:
    a = plan["plan"]["assortment"]
    ax1, ax2 = fig.subplots(1, 2)
    # SKUs before/after
    n_total = a["n_carried"] + a["n_dropped"]
    ax1.bar(["Catalogue", "Carried"], [n_total, a["n_carried"]],
            color=[GREY, BLUE], edgecolor="white")
    ax1.set_title("SKUs — before vs after", color=INK)
    ax1.set_ylabel("SKU count")
    for i, v in enumerate([n_total, a["n_carried"]]):
        ax1.text(i, v, str(v), ha="center", va="bottom", fontsize=10)
    ax1.spines[["top", "right"]].set_visible(False)
    # margin: greedy vs MILP
    ax2.bar(["Greedy", "MILP"],
            [a["greedy_annual_margin_eur"], a["optimized_annual_margin_eur"]],
            color=[PINK, GREEN], edgecolor="white")
    ax2.set_title("Assortment margin (EUR/yr)", color=INK)
    for i, v in enumerate([a["greedy_annual_margin_eur"],
                           a["optimized_annual_margin_eur"]]):
        ax2.text(i, v, _fmt_eur(v), ha="center", va="bottom", fontsize=9)
    ax2.spines[["top", "right"]].set_visible(False)
    fig.suptitle(
        f"Assortment before/after — MILP earns "
        f"{_fmt_eur(a['milp_uplift_vs_greedy_eur'])}/yr over greedy",
        fontsize=13, fontweight="bold", color=INK)


def _slide_frontier(fig, plan: dict) -> None:
    ax = fig.add_subplot(111)
    fr = _sample_frontier(plan)
    if fr is None:
        ax.text(0.5, 0.5, "Service-vs-cost frontier unavailable\n(SKU master not found)",
                ha="center", va="center", color=GREY)
        ax.axis("off")
        return
    sku, rows = fr
    x = [r["service_level"] * 100 for r in rows]
    y = [r["safety_capital_eur"] for r in rows]
    ax.plot(x, y, "-o", color=BLUE, lw=2, markerfacecolor="white")
    for xi, yi in zip(x, y, strict=False):
        ax.annotate(_fmt_eur(yi), (xi, yi), textcoords="offset points",
                    xytext=(0, 8), ha="center", fontsize=8, color=INK)
    ax.set_xlabel("Service level (%)")
    ax.set_ylabel("Safety-stock capital (EUR)")
    ax.set_title(f"Inventory service-vs-cost frontier — {sku}",
                 fontsize=14, fontweight="bold", color=INK)
    ax.spines[["top", "right"]].set_visible(False)
    ax.yaxis.grid(True, color="#eef1f5")
    ax.set_axisbelow(True)


def _slide_price_moves(fig, plan: dict) -> None:
    ax = fig.add_subplot(111)
    changes = [li["price_change_pct"] for li in _carried_lines(plan)]
    if changes:
        ax.hist(changes, bins=15, color=BLUE, edgecolor="white")
    ax.axvline(0, color=INK, lw=1)
    ax.set_xlabel("Recommended price change (%)")
    ax.set_ylabel("SKU count")
    ax.set_title("Price-move distribution (within guardrail band)",
                 fontsize=14, fontweight="bold", color=INK)
    ax.spines[["top", "right"]].set_visible(False)
    ax.yaxis.grid(True, color="#eef1f5")
    ax.set_axisbelow(True)


def _slide_promo(fig, plan: dict) -> None:
    ax = fig.add_subplot(111)
    pm = plan["plan"]["promo"]
    items = sorted(pm["allocation"].items(), key=lambda kv: -kv[1])
    cats = [c for c, _ in items]
    vals = [v for _, v in items]
    ax.barh(cats[::-1], vals[::-1], color=GREEN, edgecolor="white")
    for i, v in enumerate(vals[::-1]):
        ax.text(v, i, " " + _fmt_eur(v), va="center", fontsize=9, color=INK)
    ax.set_xlabel("Promo allocation (EUR)")
    ax.set_title(
        f"Promo allocation — {_fmt_eur(pm['incremental_margin_eur'])} "
        f"incremental ( +{_fmt_eur(pm['uplift_vs_even_eur'])} vs even split)",
        fontsize=12, fontweight="bold", color=INK)
    ax.spines[["top", "right"]].set_visible(False)


def _slide_model_quality(fig, plan: dict) -> None:
    mq = plan["model_quality"]
    fq, eq, rq = mq["forecast"], mq["elasticity"], mq["decline_risk"]
    ax1, ax2, ax3 = fig.subplots(1, 3)
    # forecast MASE vs naive
    ax1.bar(["Model", "Seasonal\nnaive"], [fq["mase"], fq["naive_mase"]],
            color=[GREEN, GREY], edgecolor="white")
    ax1.axhline(1.0, color=PINK, ls="--", lw=1)
    ax1.set_title("Forecast MASE", color=INK)
    for i, v in enumerate([fq["mase"], fq["naive_mase"]]):
        ax1.text(i, v, f"{v:.2f}", ha="center", va="bottom", fontsize=9)
    ax1.spines[["top", "right"]].set_visible(False)
    # elasticity recovery
    per = eq["per_category"]
    cats = list(per)
    trues = [per[c]["true_mean"] for c in cats]
    ests = [per[c]["estimated"] for c in cats]
    ax2.scatter(trues, ests, color=BLUE)
    lim = [min(trues + ests) - 0.2, max(trues + ests) + 0.2]
    ax2.plot(lim, lim, ls="--", color=GREY, lw=1)
    ax2.set_xlabel("True elasticity")
    ax2.set_ylabel("Estimated")
    ax2.set_title(f"Elasticity recovery\n(MAE {eq['mae_by_category']})", color=INK)
    ax2.spines[["top", "right"]].set_visible(False)
    # decline AUC
    ax3.bar(["ROC-AUC", "PR-AUC"], [rq["roc_auc"], rq["pr_auc"]],
            color=[BLUE, PINK], edgecolor="white")
    ax3.set_ylim(0, 1.05)
    ax3.set_title("Decline-risk quality", color=INK)
    for i, v in enumerate([rq["roc_auc"], rq["pr_auc"]]):
        ax3.text(i, v, f"{v:.2f}", ha="center", va="bottom", fontsize=9)
    ax3.spines[["top", "right"]].set_visible(False)
    fig.suptitle("Model quality (honest, held-out)", fontsize=14,
                 fontweight="bold", color=INK)


def _slide_actions(fig, plan: dict) -> None:
    fig.text(0.5, 0.92, "Recommended actions", ha="center", fontsize=16,
             fontweight="bold", color=INK)
    cards = plan["decision_cards"]
    y = 0.82
    for i, card in enumerate(cards, 1):
        wrapped = _wrap(card, 92)
        fig.text(0.07, y, f"{i}.", fontsize=10, color=BLUE, fontweight="bold")
        fig.text(0.11, y, wrapped, fontsize=9.5, color=INK, va="top")
        y -= 0.045 + 0.028 * wrapped.count("\n")
        if y < 0.06:
            break


def _wrap(text: str, width: int) -> str:
    import textwrap
    return "\n".join(textwrap.wrap(text, width))


def _write_pdf(plan: dict, outdir: Path) -> Path | None:
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib.backends.backend_pdf import PdfPages
    except Exception:
        return None

    builders = [
        ("full", _slide_title),
        ("axes", _slide_waterfall),
        ("subplots", _slide_assortment),
        ("axes", _slide_frontier),
        ("axes", _slide_price_moves),
        ("axes", _slide_promo),
        ("subplots", _slide_model_quality),
        ("full", _slide_actions),
    ]
    p = outdir / "executive_review.pdf"
    pngs: list = []
    with PdfPages(p) as pdf:
        for idx, (_kind, fn) in enumerate(builders, 1):
            fig = plt.figure(figsize=(11, 8.5))
            fn(fig, plan)
            fig.text(0.95, 0.03, f"{idx} / {len(builders)}", ha="right",
                     fontsize=8, color=GREY)
            fig.text(0.05, 0.03, "RevOps Optimizer — " + AUTHOR, fontsize=8,
                     color=GREY)
            pdf.savefig(fig)
            # render a PNG for the optional pptx build
            png = outdir / f"_slide_{idx}.png"
            fig.savefig(png, dpi=110)
            pngs.append(png)
            plt.close(fig)

    _write_pptx(plan, outdir, pngs)
    # clean up the intermediate PNGs
    for png in pngs:
        try:
            png.unlink()
        except OSError:
            pass
    return p


def _write_pptx(plan: dict, outdir: Path, pngs: list) -> Path | None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except Exception:
        return None
    prs = Presentation()
    prs.slide_width = Inches(11)
    prs.slide_height = Inches(8.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        if not Path(png).exists():
            continue
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(png), 0, 0, width=Inches(11),
                                 height=Inches(8.5))
    p = outdir / "executive_review.pptx"
    prs.save(p)
    return p


# --------------------------------------------------------------------------- #
# public API                                                                   #
# --------------------------------------------------------------------------- #
def write_deliverables(plan: dict, outdir: str | Path = "deliverables") -> dict:
    """Write every deliverable and return {name: path or None}."""
    out = Path(outdir)
    out.mkdir(parents=True, exist_ok=True)
    written = {
        "prescription.json": _write_json(plan, out),
        "actions.csv": _write_actions_csv(plan, out),
        "optimization_workbook.xlsx": _write_workbook(plan, out),
        "executive_review.pdf": _write_pdf(plan, out),
    }
    written["executive_review.pptx"] = (
        out / "executive_review.pptx"
        if (out / "executive_review.pptx").exists() else None)
    return written


if __name__ == "__main__":
    from .prescribe import prescribe
    plan = prescribe(verbose=False)
    result = write_deliverables(plan)
    for name, path in result.items():
        print(f"  {'ok ' if path else 'skip'} {name:<28} {path or '(library missing)'}")
